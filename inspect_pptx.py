import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation("BioPofile_Template.pptx")
def print_shape(shape, indent=""):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"{indent}Group: {shape.name}")
        for s in shape.shapes:
            print_shape(s, indent + "  ")
    else:
        if shape.has_text_frame:
            print(f"{indent}Shape: {shape.name} | Text: '{shape.text.replace(chr(10), ' ')}'")
        else:
            print(f"{indent}Shape: {shape.name} (type: {shape.shape_type})")

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i} ---")
    for shape in slide.shapes:
        print_shape(shape)
