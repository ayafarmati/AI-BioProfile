import json
import os
import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

def generate_bio_profile_dynamic(json_path="resultat_cv.json", template_path="BioPofile_Template.pptx", output_path="BioProfile_Generated.pptx"):
    if not os.path.exists(json_path):
        raise FileNotFoundError(f"Le fichier JSON {json_path} est introuvable.")
        
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    def flatten_data(data_obj, sep=", ", is_lang=False):
        if isinstance(data_obj, str):
            return data_obj
        elif isinstance(data_obj, list):
            flattened = []
            for item in data_obj:
                if isinstance(item, str):
                    flattened.append(item)
                elif isinstance(item, dict):
                    lang_name = item.get('langue') or item.get('nom')
                    if lang_name:
                        niv = item.get('niveau', '')
                        flattened.append(f"{lang_name} — {niv}" if niv else lang_name)
                    else:
                        flattened.extend(flatten_data(item, sep, is_lang).split(sep))
            return sep.join(flattened)
        elif isinstance(data_obj, dict):
            flattened = []
            for k, v in data_obj.items():
                if isinstance(v, str) and is_lang:
                    flattened.append(f"{k} — {v}")
                elif isinstance(v, str):
                    flattened.append(v)
                elif isinstance(v, list):
                    flattened.extend(flatten_data(v, sep, is_lang).split(sep))
                elif isinstance(v, dict):
                    flattened.extend(flatten_data(v, sep, is_lang).split(sep))
            return sep.join(flattened)
        return str(data_obj)
        
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Préparation des données formatées
    nom = data.get("nom_complet", "")
    titre = data.get("titre_professionnel", "")
    dispo = data.get("disponibilite", "Immédiate")
    
    langues = flatten_data(data.get("langues", []), sep="\n", is_lang=True)
    hard = flatten_data(data.get("hard_skills", []))
    soft = flatten_data(data.get("soft_skills", []))
    outils = flatten_data(data.get("outils_et_technologies", []))
    
    projets = data.get("projets_et_experiences", [])
    if isinstance(projets, dict):
        for v in projets.values():
            if isinstance(v, list):
                projets = v
    hard = flatten_data(data.get("hard_skills", []))
    soft = flatten_data(data.get("soft_skills", []))
    outils = flatten_data(data.get("outils_et_technologies", []))

    # Fixed base size for skills, let PowerPoint AutoFit handle any shrinking
    uniform_skill_size = 14

    # Dictionnaire de remplacement simple
    replacements = {
        "{{NOM}}": nom,
        "{{TITRE}}": titre,
        "{{DISPO}}": dispo,
        "{{LANGUES}}": langues,
        "{{HARD}}": hard,
        "{{SOFT}}": soft,
        "{{OUTILS}}": outils
    }

    # Liste pour mémoriser les formes à supprimer (comme {{PHOTO}})
    shapes_to_delete = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        text = shape.text
        
        # 1. Gestion de l'Image de profil
        if "{{PHOTO}}" in text:
            # Récupérer les coordonnées et dimensions
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            
            photo_path = data.get("photo_path")
            if photo_path:
                if "?" in photo_path:
                    photo_path = photo_path.split("?")[0]
                if photo_path.startswith("/"):
                    photo_path = photo_path[1:]
                
            if photo_path and os.path.exists(photo_path):
                try:
                    slide.shapes.add_picture(photo_path, left, top, width=width, height=height)
                except Exception as e:
                    print(f"Erreur insertion image: {e}")
                    
            # Marquer la boîte de texte {{PHOTO}} pour suppression
            shapes_to_delete.append(shape)
            continue

        # 2. Gestion des Projets (Formatage complexe avec Titre et Description)
        if "{{PROJETS}}" in text:
            # Sauvegarder le style initial de la balise
            font_size = None
            font_bold = None
            font_color = None
            font_name = None
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                r = shape.text_frame.paragraphs[0].runs[0]
                font_size = r.font.size
                font_bold = r.font.bold
                if hasattr(r.font.color, 'rgb'):
                    font_color = r.font.color.rgb
                font_name = r.font.name

            shape.text_frame.clear()
            shape.text_frame.word_wrap = True
            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            if isinstance(projets, list):
                # Always use readable sizes; if there are many projects, let PowerPoint's AutoFit (TEXT_TO_FIT_SHAPE) scale them down naturally.
                base_title_size = 18
                desc_size = 14
                space_size = 12
                
                if len(projets) == 4:
                    base_title_size = 16
                    desc_size = 13
                    space_size = 8
                elif len(projets) >= 5:
                    base_title_size = 14
                    desc_size = 12
                    space_size = 6
                
                for i, p_data in enumerate(projets):
                    titre_proj = p_data.get("titre", "") if isinstance(p_data, dict) else ""
                    desc_proj = p_data.get("description", "") if isinstance(p_data, dict) else str(p_data)
                    
                    # Ajouter le titre (Hot Pink)
                    from pptx.dml.color import RGBColor
                    p = shape.text_frame.add_paragraph()
                    p.text = titre_proj
                    if font_name: p.font.name = font_name
                    p.font.size = Pt(base_title_size)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 105, 180) # Hot Pink comme l'original
                    
                    # Ajouter la description (Noir et Bold)
                    p = shape.text_frame.add_paragraph()
                    p.text = desc_proj
                    if font_name: p.font.name = font_name
                    p.font.size = Pt(desc_size)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # Espace entre les projets
                    if i < len(projets) - 1:
                        p_space = shape.text_frame.add_paragraph()
                        p_space.font.size = Pt(space_size)
            continue

        # 3. Gestion des remplacements classiques
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                for tag, value in replacements.items():
                    if tag in r.text:
                        str_val = str(value)
                        r.text = r.text.replace(tag, str_val)
                        shape.text_frame.word_wrap = True
                        
                        # Application de la taille uniforme pour les compétences
                        if tag in ["{{HARD}}", "{{SOFT}}", "{{OUTILS}}"]:
                            r.font.size = Pt(uniform_skill_size)
                        else:
                            # Auto-ajustement manuel basé sur la longueur du texte pour les autres tags
                            if len(str_val) > 70 and r.font.size:
                                reduction = int((len(str_val) - 70) / 40)
                                new_size = max(11, r.font.size.pt - reduction)
                                r.font.size = Pt(new_size)

    # Supprimer les formes obsolètes (ex: boîte {{PHOTO}})
    for shape in shapes_to_delete:
        sp = shape.element
        sp.getparent().remove(sp)

    prs.save(output_path)
    print(f"Presentation dynamique generee: {output_path}")
    return output_path

if __name__ == "__main__":
    generate_bio_profile_dynamic()
