import json
import os
import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

def generate_bio_profile_dynamic(json_path="resultat_cv.json", template_path="BioProfile_OFF.pptx", output_path="BioProfile_Generated.pptx"):
    if not os.path.exists(json_path):
        raise FileNotFoundError(f"Le fichier JSON {json_path} est introuvable.")
        
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    def flatten_data(data_obj, sep=", ", is_lang=False):
        if isinstance(data_obj, str):
            return data_obj
        elif isinstance(data_obj, list):
            flattened = []
            for item in data_obj:
                if isinstance(item, str):
                    flattened.append(item)
                elif isinstance(item, dict):
                    lang_name = item.get('langue') or item.get('nom')
                    if lang_name:
                        niv = item.get('niveau', '')
                        flattened.append(f"{lang_name} — {niv}" if niv else lang_name)
                    else:
                        flattened.extend(flatten_data(item, sep, is_lang).split(sep))
            return sep.join(flattened)
        elif isinstance(data_obj, dict):
            flattened = []
            for k, v in data_obj.items():
                if isinstance(v, str) and is_lang:
                    flattened.append(f"{k} — {v}")
                elif isinstance(v, str):
                    flattened.append(v)
                elif isinstance(v, list):
                    flattened.extend(flatten_data(v, sep, is_lang).split(sep))
                elif isinstance(v, dict):
                    flattened.extend(flatten_data(v, sep, is_lang).split(sep))
            return sep.join(flattened)
        return str(data_obj)
        
    def format_as_newlines(text):
        if not text: return ""
        # Remove any fake bullets from LLM
        text = text.replace('•', '').replace('-', '')
        # Split by newline or comma
        text = text.replace('\n', ',')
        items = [x.strip() for x in text.split(",") if x.strip()]
        return "\n".join(items)

    nom = data.get("nom_complet", "")
    titre = data.get("titre_professionnel", "")
    dispo = data.get("disponibilite", "Immédiate")
    if len(dispo) > 25:
        dispo = "Immédiate"
        
    langues = format_as_newlines(flatten_data(data.get("langues", []), sep="\n", is_lang=True))
    hard = format_as_newlines(flatten_data(data.get("hard_skills", [])))
    soft = format_as_newlines(flatten_data(data.get("soft_skills", [])))
    outils = format_as_newlines(flatten_data(data.get("outils_et_technologies", [])))
    
    projets = data.get("projets_et_experiences", [])
    if isinstance(projets, dict):
        for v in projets.values():
            if isinstance(v, list):
                projets = v
                
    # Formatage des projets pour la zone blanche (Max 3)
    descs_projets_list = []
    if isinstance(projets, list):
        for p in projets[:3]:
            t = p.get("titre", "").strip()
            d = p.get("description", "").strip()
            
            # Nettoyage des puces IA
            t = t.replace('•', '').replace('-', '').strip()
            d = d.replace('•', '').replace('-', '').strip()
            
            # On met tout sur une seule ligne (ou quelques lignes)
            if d:
                descs_projets_list.append(f"{t} : {d}")
            elif t:
                descs_projets_list.append(t)
                
    descs_projets = "\n".join(descs_projets_list)

    replacements = {
        "{{NOM}}": nom,
        "{{TITRE}}": titre,
        "{{DISPO}}": dispo,
        "{{LANGUES}}": langues,
        "{{HARD}}": hard,
        "{{SOFT}}": soft,
        "{{OUTILS}}": outils,
        "{{PROJET_TITRE}}": "", # Laissé vide exprès pour éviter le dépassement dans la bannière bleue
        "{{PROJET_DESC}}": descs_projets
    }

    def get_all_shapes(shapes):
        for shape in shapes:
            if shape.has_text_frame:
                yield shape
            if shape.shape_type == 6: # Group shape
                yield from get_all_shapes(shape.shapes)

    def replace_tag_with_paragraphs(shape, tag, replacement_text):
        target_p = None
        target_r = None
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if tag in r.text:
                    target_p = p
                    target_r = r
                    break
            if target_p:
                break
                
        if not target_p:
            return False
            
        lines = str(replacement_text).split('\n')
        
        # Replace the text in the first run
        target_r.text = target_r.text.replace(tag, lines[0])
        
        # For the remaining lines, add new paragraphs
        for line in lines[1:]:
            new_p = shape.text_frame.add_paragraph()
            new_run = new_p.add_run()
            new_run.text = line
            if target_r.font.size:
                new_run.font.size = target_r.font.size
            if target_r.font.bold is not None:
                new_run.font.bold = target_r.font.bold
            new_p.level = target_p.level # Conserve la puce native PPTX
            
        return True

    prs = Presentation(template_path)
    slide = prs.slides[0]
    
    shapes_to_delete = []

    for shape in get_all_shapes(slide.shapes):
        text = shape.text
        
        # 1. Gestion de l'Image de profil
        if "{{PHOTO}}" in text:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            
            photo_path = data.get("photo_path")
            if photo_path:
                if "?" in photo_path:
                    photo_path = photo_path.split("?")[0]
                if photo_path.startswith("/"):
                    photo_path = photo_path[1:]
                
            if photo_path and os.path.exists(photo_path):
                try:
                    slide.shapes.add_picture(photo_path, left, top, width=width, height=height)
                except Exception as e:
                    print(f"Erreur insertion image: {e}")
                    
            shapes_to_delete.append(shape)
            continue

        # 2. Remplacements intelligents avec préservation des puces
        if shape.has_text_frame:
            for tag, value in replacements.items():
                if tag in shape.text:
                    replace_tag_with_paragraphs(shape, tag, value)
                    shape.text_frame.word_wrap = True
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    
                    # Réduction dynamique de la taille de police pour le TITRE ou NOM s'il est trop long
                    if tag in ["{{TITRE}}", "{{NOM}}"] and len(str(value)) > 40:
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                if r.font.size:
                                    reduction = int((len(str(value)) - 40) / 15)
                                    new_size = max(10, r.font.size.pt - reduction)
                                    r.font.size = Pt(new_size)

            # Nettoyage des paragraphes vides de fin (ex: \n\n\n\n laissés par le template)
            p_elements = [p._element for p in shape.text_frame.paragraphs]
            for p_elem in reversed(p_elements):
                if not p_elem.text.strip():
                    try:
                        pass # p_elem.getparent().remove(p_elem)
                    except:
                        pass
                else:
                    break

    for shape in shapes_to_delete:
        sp = shape.element
        sp.getparent().remove(sp)

    prs.save(output_path)
    print(f"Presentation dynamique generee: {output_path}")
    return output_path

if __name__ == "__main__":
    generate_bio_profile_dynamic()
