import json
import os
import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

def generate_bio_profile_dynamic(json_path="resultat_cv.json", template_path="BioProfile_OFF.pptx", output_path="BioProfile_Generated.pptx"):
    if not os.path.exists(json_path):
        raise FileNotFoundError(f"Le fichier JSON {json_path} est introuvable.")
        
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    def flatten_data(data_obj, sep=", ", is_lang=False):
        if isinstance(data_obj, str):
            return data_obj
        elif isinstance(data_obj, list):
            flattened = []
            for item in data_obj:
                if isinstance(item, str):
                    flattened.append(item)
                elif isinstance(item, dict):
                    lang_name = item.get('langue') or item.get('nom')
                    if lang_name:
                        niv = item.get('niveau', '')
                        flattened.append(f"{lang_name} — {niv}" if niv else lang_name)
                    else:
                        flattened.extend(flatten_data(item, sep, is_lang).split(sep))
            return sep.join(flattened)
        elif isinstance(data_obj, dict):
            flattened = []
            for k, v in data_obj.items():
                if isinstance(v, str) and is_lang:
                    flattened.append(f"{k} — {v}")
                elif isinstance(v, str):
                    flattened.append(v)
                elif isinstance(v, list):
                    flattened.extend(flatten_data(v, sep, is_lang).split(sep))
                elif isinstance(v, dict):
                    flattened.extend(flatten_data(v, sep, is_lang).split(sep))
            return sep.join(flattened)
        return str(data_obj)
        
    def format_as_newlines(text):
        if not text: return ""
        import re
        if '\n' in text:
            # It's already multiline
            lines = [line.strip().lstrip('•').lstrip('-').strip() for line in text.split('\n')]
            return "\n".join([x for x in lines if x])
        else:
            # Try splitting by '. ' for sentences/categories
            if '. ' in text:
                items = [x.strip() for x in text.split('. ') if x.strip()]
                return "\n".join(items)
            # Otherwise split by comma not inside parens
            items = [x.strip() for x in re.split(r',\s*(?![^()]*\))', text) if x.strip()]
            return "\n".join(items)

    nom = data.get("nom_complet", "")
    titre = data.get("titre_professionnel", "")
    dispo = data.get("disponibilite", "Immédiate")
    if len(dispo) > 25:
        dispo = "Immédiate"
        
    langues = format_as_newlines(flatten_data(data.get("langues", []), sep="\n", is_lang=True))
    hard = format_as_newlines(flatten_data(data.get("hard_skills", [])))
    soft = format_as_newlines(flatten_data(data.get("soft_skills", [])))
    outils = format_as_newlines(flatten_data(data.get("outils_et_technologies", [])))
    
    projets = data.get("projets_et_experiences", [])
    if isinstance(projets, dict):
        for v in projets.values():
            if isinstance(v, list):
                projets = v
                
    # Formatage des projets
    descs_projets_list = []
    titres_projets_list = []
    if isinstance(projets, list):
        for p in projets[:3]:
            t = p.get("titre", "").strip()
            d = p.get("description", "").strip()
            
            # Nettoyage des puces IA pour le titre
            t = t.replace('•', '').replace('-', '').strip()
            if t:
                titres_projets_list.append(t)
            
            if d:
                # Formatage de la description avec des puces intelligentes
                import re
                if '\n' in d:
                    lines = [line.strip().lstrip('•').lstrip('-').strip() for line in d.split('\n')]
                    d_formatted = "\n".join([x for x in lines if x])
                else:
                    # Splitting by '. ' to create bullet points for experiences
                    items = [x.strip() for x in d.split('. ') if x.strip()]
                    d_formatted = "\n".join(items)
                
                descs_projets_list.append(d_formatted)
                
    descs_projets = "\n\n".join(descs_projets_list)
    titres_projets = "\n\n".join(titres_projets_list)

    replacements = {
        "{{NOM}}": nom,
        "{{TITRE}}": titre,
        "{{DISPO}}": dispo,
        "{{LANGUES}}": langues,
        "{{HARD}}": hard,
        "{{SOFT}}": soft,
        "{{OUTILS}}": outils,
        "{{PROJET_TITRE}}": titres_projets,
        "{{PROJET_DESC}}": descs_projets
    }

    def get_all_shapes(shapes):
        for shape in shapes:
            if shape.has_text_frame:
                yield shape
            if shape.shape_type == 6: # Group shape
                yield from get_all_shapes(shape.shapes)

    def replace_tag_with_paragraphs(shape, tag, replacement_text):
        target_p = None
        target_r = None
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if tag in r.text:
                    target_p = p
                    target_r = r
                    break
            if target_p:
                break
                
        if not target_p:
            return False
            
        lines = str(replacement_text).split('\n')
        
        # Calculate dynamic font size to prevent overflow
        if target_r.font.size:
            original_size_pt = target_r.font.size.pt
        else:
            original_size_pt = 14 # default

        new_size_pt = original_size_pt
        
        if tag in ["{{TITRE}}", "{{NOM}}"] and len(str(replacement_text)) > 40:
            reduction = int((len(str(replacement_text)) - 40) / 15)
            new_size_pt = max(10, original_size_pt - reduction)
        elif len(lines) > 6 and tag not in ["{{PROJET_DESC}}"]:
            # Reduce font size significantly for long lists (e.g., Technical Skills)
            reduction = (len(lines) - 6) * 1.0
            new_size_pt = max(8, original_size_pt - reduction)

        # Apply new size to the first run
        target_r.font.size = Pt(new_size_pt)
        target_r.text = target_r.text.replace(tag, lines[0])
        
        # If font size was reduced, we must also remove fixed line spacing from the paragraph 
        # to prevent huge vertical gaps between tiny text.
        if new_size_pt < original_size_pt and target_p._element.pPr is not None:
            for child in list(target_p._element.pPr):
                if child.tag.endswith('lnSpc') or child.tag.endswith('spcAft') or child.tag.endswith('spcBef'):
                    target_p._element.pPr.remove(child)
        
        current_p_elem = target_p._element
        
        # For the remaining lines, add new paragraphs and move them to the correct position
        for line in lines[1:]:
            new_p = shape.text_frame.add_paragraph()
            new_run = new_p.add_run()
            
            # Copy all run properties to preserve color and styling
            if hasattr(target_r, '_r') and target_r._r.rPr is not None:
                new_rPr = copy.deepcopy(target_r._r.rPr)
                # Remove the size from copied rPr so we can set our own dynamically calculated size
                for child in list(new_rPr):
                    if child.tag.endswith('sz'):
                        new_rPr.remove(child)
                new_run._r.insert(0, new_rPr)
                
            new_run.text = line
            new_run.font.size = Pt(new_size_pt)
            
            if target_r.font.bold is not None:
                new_run.font.bold = target_r.font.bold
                
            # Deep copy paragraph properties to preserve native PowerPoint bullets exactly!
            if target_p._element.pPr is not None:
                new_pPr = copy.deepcopy(target_p._element.pPr)
                new_p._element.insert(0, new_pPr)
            
            # Move it right after the current paragraph to maintain exact order
            current_p_elem.addnext(new_p._element)
            current_p_elem = new_p._element
            
        return True

    prs = Presentation(template_path)
    slide = prs.slides[0]
    
    shapes_to_delete = []

    for shape in get_all_shapes(slide.shapes):
        text = shape.text
        
        # 1. Gestion de l'Image de profil
        if "{{PHOTO}}" in text:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            
            photo_path = data.get("photo_path")
            if photo_path:
                if "?" in photo_path:
                    photo_path = photo_path.split("?")[0]
                if photo_path.startswith("/"):
                    photo_path = photo_path[1:]
                
            if photo_path and os.path.exists(photo_path):
                try:
                    slide.shapes.add_picture(photo_path, left, top, width=width, height=height)
                except Exception as e:
                    print(f"Erreur insertion image: {e}")
                    
            shapes_to_delete.append(shape)
            continue

        # 2. Remplacements intelligents avec préservation des puces
        if shape.has_text_frame:
            for tag, value in replacements.items():
                if tag in shape.text:
                    replace_tag_with_paragraphs(shape, tag, value)
                    shape.text_frame.word_wrap = True
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            # Nettoyage sécurisé des paragraphes vides de fin (ex: \n\n\n\n laissés par le template)
            p_elements = [p._element for p in shape.text_frame.paragraphs]
            # On garde toujours au moins 1 paragraphe pour ne pas corrompre le fichier XML PPTX
            while len(p_elements) > 1:
                p_elem = p_elements[-1]
                if not p_elem.text.strip():
                    try:
                        p_elem.getparent().remove(p_elem)
                        p_elements.pop()
                    except:
                        break
                else:
                    break

    for shape in shapes_to_delete:
        sp = shape.element
        sp.getparent().remove(sp)

    prs.save(output_path)
    print(f"Presentation dynamique generee: {output_path}")
    return output_path

if __name__ == "__main__":
    generate_bio_profile_dynamic()
