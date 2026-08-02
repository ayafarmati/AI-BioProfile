import json
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def generate_bio_profile(json_path="resultat_cv.json", template_path="BioPofile_Template.pptx", output_path="BioProfile_Generated.pptx"):
    if not os.path.exists(json_path):
        raise FileNotFoundError(f"Le fichier JSON {json_path} est introuvable.")
        
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    def flatten_data(data_obj, sep=", ", is_lang=False):
        if isinstance(data_obj, str):
            return data_obj
        elif isinstance(data_obj, list):
            flattened = []
            for item in data_obj:
                if isinstance(item, str):
                    flattened.append(item)
                elif isinstance(item, dict):
                    if 'langue' in item:
                        niv = item.get('niveau', '')
                        flattened.append(f"{item['langue']} - {niv}" if niv else item['langue'])
                    else:
                        flattened.extend(flatten_data(item, sep, is_lang).split(sep))
            return sep.join(flattened)
        elif isinstance(data_obj, dict):
            flattened = []
            for k, v in data_obj.items():
                if isinstance(v, str) and is_lang:
                    flattened.append(f"{k} - {v}")
                elif isinstance(v, str):
                    flattened.append(v)
                elif isinstance(v, list):
                    flattened.extend(flatten_data(v, sep, is_lang).split(sep))
                elif isinstance(v, dict):
                    flattened.extend(flatten_data(v, sep, is_lang).split(sep))
            return sep.join(flattened)
        return str(data_obj)
        
    prs = Presentation(template_path)
    slide = prs.slides[0]
    
    def add_text_box(left, top, width, height, text, font_size=20, bold=False, alignment=None, rgb=RGBColor(0, 0, 0)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = rgb
        if alignment:
            p.alignment = alignment
        return txBox

    # 1. Image de profil
    photo_path = data.get("photo_path")
    if photo_path and photo_path.startswith("/"):
        photo_path = photo_path[1:]
        
    if photo_path and os.path.exists(photo_path):
        try:
            slide.shapes.add_picture(photo_path, Emu(1085817), Emu(1314776), width=Emu(3041855), height=Emu(2754260))
        except Exception as e:
            print(f"Erreur image: {e}")

    # Nettoyer et remplacer les textes existants pour eviter la superposition
    nom = data.get("nom_complet", "")
    dispo_text = data.get("disponibilite", "Immédiate")
    
    nom_top = Emu(4240485)
    nom_left = Emu(1180138)
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_upper = shape.text.upper()
        
        # Remplacer "Nom Complet"
        if "NOM COMPLET" in text_upper:
            nom_top = shape.top
            nom_left = shape.left
            shape.text_frame.clear()
            p = shape.text_frame.add_paragraph()
            p.text = nom
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255) # Blanc sur fond bleu
            p.alignment = PP_ALIGN.CENTER
            nom = "" # Marquer comme fait
            
        # Remplacer "DISPONIBILITE"
        elif "DISPONIBILIT" in text_upper:
            shape.text_frame.clear()
            p = shape.text_frame.add_paragraph()
            
            run1 = p.add_run()
            run1.text = "DISPONIBILITÉ:   "
            run1.font.size = Pt(26)
            run1.font.bold = True
            run1.font.color.rgb = RGBColor(0xFF, 0x52, 0xBA)
            
            run2 = p.add_run()
            run2.text = dispo_text
            run2.font.size = Pt(26)
            run2.font.bold = True
            run2.font.color.rgb = RGBColor(255, 255, 255)
            
    # Si le champ Nom Complet n'existait pas, on l'ajoute
    if nom:
        add_text_box(Emu(1180138), Emu(4240485), Emu(2853214), Emu(562200), nom, font_size=30, bold=True, alignment=PP_ALIGN.CENTER, rgb=RGBColor(255, 255, 255))
        
    titre = data.get("titre_professionnel", "")
    if titre:
        add_text_box(nom_left, nom_top + Emu(600000), Emu(3000000), Emu(500000), titre, font_size=20, bold=False, alignment=PP_ALIGN.CENTER, rgb=RGBColor(200, 200, 200))

    # 4. Langues
    langues = data.get("langues", [])
    langues_str = flatten_data(langues, sep="\n", is_lang=True)
    add_text_box(Emu(588266), Emu(7000000), Emu(4036957), Emu(1785964), langues_str, font_size=20, bold=True, rgb=RGBColor(255, 255, 255))

    # 5. Projets - UTILISATION D'UNE SEULE BOX POUR EVITER LE CHEVAUCHEMENT
    projets = data.get("projets_et_experiences", [])
    if isinstance(projets, dict):
        for v in projets.values():
            if isinstance(v, list):
                projets = v
                break

    if projets:
        txBox_proj = slide.shapes.add_textbox(Emu(6388294), Emu(1000000), Emu(11297755), Emu(4000000))
        tf_proj = txBox_proj.text_frame
        tf_proj.word_wrap = True
        tf_proj.clear()
        
        if isinstance(projets, list):
            for i, p in enumerate(projets[:3]):
                titre_proj = p.get("titre", "") if isinstance(p, dict) else ""
                desc_proj = p.get("description", "") if isinstance(p, dict) else str(p)
                
                # Titre en Gras
                p_title = tf_proj.add_paragraph()
                p_title.text = titre_proj
                p_title.font.size = Pt(22)
                p_title.font.bold = True
                p_title.font.color.rgb = RGBColor(0, 0, 0)
                
                # Description
                p_desc = tf_proj.add_paragraph()
                p_desc.text = desc_proj
                p_desc.font.size = Pt(18)
                p_desc.font.bold = False
                p_desc.font.color.rgb = RGBColor(0, 0, 0)
                
                # Espace
                p_space = tf_proj.add_paragraph()
                p_space.font.size = Pt(14)

    # 6. Hard Skills
    skills = data.get("hard_skills", [])
    skills_text = flatten_data(skills)
    add_text_box(Emu(5821465), Emu(5815490), Emu(3115572), Emu(4132468), skills_text, font_size=18, bold=True)

    # 7. Soft Skills
    soft_skills = data.get("soft_skills", [])
    soft_skills_text = flatten_data(soft_skills)
    add_text_box(Emu(9546637), Emu(6020278), Emu(3565954), Emu(3722893), soft_skills_text, font_size=18, bold=True)
    
    # 8. Outils & Technologies
    outils = data.get("outils_et_technologies", [])
    outils_text = flatten_data(outils)
    add_text_box(Emu(13655839), Emu(5944982), Emu(3962469), Emu(3313318), outils_text, font_size=18, bold=True)

    prs.save(output_path)
    print(f"Presentation generee: {output_path}")
    return output_path

if __name__ == "__main__":
    generate_bio_profile()
